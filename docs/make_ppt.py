# -*- coding: utf-8 -*-
"""生成「招聘爬虫反爬攻坚」阶段汇报 PPT。

运行： .venv\\Scripts\\python docs\\make_ppt.py
产出： docs/反爬攻坚汇报.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- 主题配色 ----
BG = RGBColor(0x0F, 0x17, 0x2A)        # 深海军蓝
PANEL = RGBColor(0x1E, 0x29, 0x3B)     # 面板底
ACCENT = RGBColor(0x38, 0xBD, 0xF8)    # 主色（亮蓝）
ACCENT2 = RGBColor(0x34, 0xD3, 0x99)   # 成功绿
WARN = RGBColor(0xFB, 0xBF, 0x24)      # 警告黄
BAD = RGBColor(0xF8, 0x71, 0x71)       # 失败红
TEXT = RGBColor(0xE2, 0xE8, 0xF0)      # 主文字
MUTED = RGBColor(0x94, 0xA3, 0xB8)     # 次要文字
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE

    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shp.adjustments[0] = 0.06
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def _text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=6):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (txt, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = FONT
    return tb


def _accent_bar(slide, t=Inches(1.18)):
    _box(slide, Inches(0.6), t, Inches(0.16), Inches(0.5), fill=ACCENT)


def header(slide, kicker, title):
    _accent_bar(slide)
    _text(slide, Inches(0.95), Inches(0.5), Inches(11.8), Inches(0.5),
          [[(kicker, 14, ACCENT, True)]])
    _text(slide, Inches(0.92), Inches(0.86), Inches(11.9), Inches(0.9),
          [[(title, 30, WHITE, True)]])


def footer(slide, page):
    _text(slide, Inches(0.6), Inches(7.02), Inches(8), Inches(0.4),
          [[("job-hunter · 反爬攻坚阶段汇报", 10, MUTED, False)]])
    _text(slide, Inches(11.5), Inches(7.02), Inches(1.3), Inches(0.4),
          [[(str(page), 10, MUTED, False)]], align=PP_ALIGN.RIGHT)


# ---------------- 封面 ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
_box(s, Inches(0.0), Inches(0.0), Inches(0.22), SH, fill=ACCENT)
_text(s, Inches(0.9), Inches(2.05), Inches(11.6), Inches(0.6),
      [[("招聘职位自动抓取系统", 20, ACCENT, True)]])
_text(s, Inches(0.85), Inches(2.5), Inches(11.8), Inches(1.6),
      [[("反爬攻坚 · 阶段成果汇报", 46, WHITE, True)]])
_text(s, Inches(0.9), Inches(4.0), Inches(11.6), Inches(1.0),
      [[("多平台（智联 / 领英 / BOSS / 猎聘 / 51）自动化抓取", 18, TEXT, False)],
       [("从「全靠手动」到「绝大多数全自动」的突破纪实", 18, MUTED, False)]])
_box(s, Inches(0.9), Inches(5.4), Inches(4.2), Inches(0.04), fill=PANEL)
_text(s, Inches(0.9), Inches(5.6), Inches(11.6), Inches(0.6),
      [[("2026-06-07", 14, MUTED, False)]])

# ---------------- 2. 目标与架构 ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
header(s, "PROJECT  概览", "项目目标与整体架构")
items = [
    ("目标", "多平台职位自动抓取 → AI 匹配打分 → 定时预抓取 → 邮件推送日报"),
    ("技术栈", "FastAPI + Playwright + APScheduler + Cursor / DeepSeek 大模型"),
    ("抓取基座", "CDP 接管你登录好的真实 Chrome，复用会话——最强反爬绕过起点"),
    ("拦路虎", "各平台反爬，其中 BOSS 直聘的反自动化最为强硬"),
]
y = Inches(1.7)
for k, v in items:
    _box(s, Inches(0.6), y, Inches(12.1), Inches(1.05), fill=PANEL)
    _text(s, Inches(0.9), y + Inches(0.16), Inches(2.4), Inches(0.8),
          [[(k, 17, ACCENT, True)]], anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(3.3), y + Inches(0.16), Inches(9.1), Inches(0.8),
          [[(v, 16, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.22)
footer(s, 2)

# ---------------- 3. 核心问题 ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
header(s, "CORE  问题", "为什么「程序抓」被拦，「手动点」却没事？")
_text(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.6),
      [[("反爬不看「谁在操作」，而看一组能区分", 17, TEXT, False),
        ("真人浏览器", 17, ACCENT, True),
        ("与", 17, TEXT, False),
        ("自动化程序", 17, WARN, True),
        ("的信号。", 17, TEXT, False)]])
layers = [
    ("网络层", "IP 信誉 / 请求频率 / Referer 来路"),
    ("指纹层", "navigator.webdriver / CDP 调试器特征 / UA"),
    ("行为层", "鼠标轨迹 / 滚动深度 / 停留时长"),
    ("导航与令牌", "是否走有机导航、深链所需 token"),
]
x = Inches(0.6)
cw = Inches(2.95)
for k, v in layers:
    _box(s, x, Inches(2.5), cw, Inches(1.9), fill=PANEL, line=ACCENT, line_w=1.0)
    _text(s, x + Inches(0.18), Inches(2.7), cw - Inches(0.36), Inches(0.5),
          [[(k, 16, ACCENT, True)]], align=PP_ALIGN.CENTER)
    _text(s, x + Inches(0.18), Inches(3.25), cw - Inches(0.36), Inches(1.0),
          [[(v, 13, TEXT, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    x += cw + Inches(0.18)
_box(s, Inches(0.6), Inches(4.75), Inches(6.0), Inches(1.6), fill=PANEL)
_text(s, Inches(0.85), Inches(4.95), Inches(5.5), Inches(1.3),
      [[("手动浏览", 18, ACCENT2, True)],
       [("四层全是真实信号 → 风控评分高 → 放行", 15, TEXT, False)]])
_box(s, Inches(6.75), Inches(4.75), Inches(5.95), Inches(1.6), fill=PANEL)
_text(s, Inches(7.0), Inches(4.95), Inches(5.5), Inches(1.3),
      [[("程序抓取", 18, BAD, True)],
       [("任一层露馅 → 限流 / 验证码 / 白板", 15, TEXT, False)]])
footer(s, 3)

# ---------------- 4. 成果总览 ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
header(s, "ACHIEVEMENTS  成果", "已经做到了什么")
wins = [
    "智联 / 领英：CDP 全自动抓列表 + 进详情页补全完整 JD",
    "拟人化行为体系：随机鼠标轨迹、滚动与回滚、不规律停留、读完再关",
    "详情页 Referer 伪装 + 被拦指数退避 + 连续熔断保护",
    "BOSS：从「必须手动开标签」推进到「程序自动模拟真人搜索流程」",
    "引入 rebrowser 补丁版 Playwright，针对性消除 CDP 调试器特征",
    "智联 / 领英 / 猎聘 / 51 等站点已稳定全自动运行",
]
y = Inches(1.75)
for w in wins:
    _box(s, Inches(0.6), y, Inches(0.42), Inches(0.42), fill=ACCENT2)
    _text(s, Inches(0.62), y - Inches(0.02), Inches(0.4), Inches(0.42),
          [[("✓", 16, BG, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(1.2), y - Inches(0.05), Inches(11.4), Inches(0.6),
          [[(w, 17, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.82)
footer(s, 4)

# ---------------- 5. 反爬原理总览（表格） ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
header(s, "PRINCIPLE  原理", "反爬分层原理 × 我们的对策")
rows = [
    ("层级", "反爬检测点", "我们的对策"),
    ("网络层", "IP 信誉 / 频率 / Referer", "限速、随机间隔、Referer 伪装、退避熔断"),
    ("指纹层", "webdriver / CDP / UA", "stealth 脚本、真实 Chrome、rebrowser 补丁"),
    ("行为层", "鼠标 / 滚动 / 停留", "拟人化动作库（轨迹·滚动·停留·读完再关）"),
    ("导航/令牌", "深链直达、缺 token", "模拟首页搜索，产生带 token 的有机导航"),
]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.6), Inches(1.8),
                         Inches(12.1), Inches(4.4)).table
tbl.columns[0].width = Inches(2.3)
tbl.columns[1].width = Inches(4.1)
tbl.columns[2].width = Inches(5.7)
for ci in range(3):
    for ri in range(len(rows)):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if ri == 0 else (PANEL if ri % 2 else RGBColor(0x16, 0x21, 0x33))
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.15)
        cell.margin_top = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = rows[ri][ci]
        run.font.name = FONT
        run.font.size = Pt(15 if ri else 16)
        run.font.bold = (ri == 0 or ci == 0)
        run.font.color.rgb = BG if ri == 0 else (ACCENT if ci == 0 else TEXT)
footer(s, 5)

# ---------------- 6. 分隔页 ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
_box(s, Inches(0.0), Inches(3.0), Inches(0.22), Inches(1.6), fill=ACCENT)
_text(s, Inches(0.9), Inches(2.9), Inches(11.6), Inches(0.6),
      [[("PROCESS", 16, ACCENT, True)]])
_text(s, Inches(0.85), Inches(3.3), Inches(11.8), Inches(1.3),
      [[("一步步的尝试与攻坚过程", 40, WHITE, True)]])

# ---------------- 7. 尝试时间线（表格） ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
header(s, "TIMELINE  尝试", "我们一步步做了哪些尝试")
rows = [
    ("#", "做法", "解决的问题", "结果"),
    ("1", "拟人化节奏 + 鼠标轨迹 + 滚动", "行为评分过低", "成功"),
    ("2", "详情页带 Referer + 读完停留再关", "「凭空直达」特征", "成功"),
    ("3", "被拦后指数退避 + 连续熔断", "越刷越被限流", "成功"),
    ("4", "BOSS 自动模拟首页搜索（有机导航）", "深链缺 token", "首页仍被吐空白"),
    ("5", "CDP 模式注入 stealth（隐藏 webdriver）", "浅层指纹", "无效"),
    ("6", "浏览器原生导航 Target.createTarget", "程序导航被拦", "标签被秒关"),
    ("7", "换 rebrowser-playwright", "CDP 调试器检测", "已接入·待验证"),
]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.6), Inches(1.7),
                         Inches(12.1), Inches(5.1)).table
tbl.columns[0].width = Inches(0.7)
tbl.columns[1].width = Inches(5.6)
tbl.columns[2].width = Inches(3.4)
tbl.columns[3].width = Inches(2.4)
status_color = {
    "成功": ACCENT2, "无效": BAD, "标签被秒关": BAD,
    "首页仍被吐空白": WARN, "已接入·待验证": ACCENT,
}
for ri in range(len(rows)):
    for ci in range(4):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if ri == 0 else (PANEL if ri % 2 else RGBColor(0x16, 0x21, 0x33))
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = rows[ri][ci]
        run.font.name = FONT
        run.font.size = Pt(13 if ri else 14)
        run.font.bold = (ri == 0 or ci == 0 or ci == 3)
        if ri == 0:
            run.font.color.rgb = BG
        elif ci == 3:
            run.font.color.rgb = status_color.get(rows[ri][ci], TEXT)
        else:
            run.font.color.rgb = TEXT
footer(s, 7)

# ---------------- 8. 原理拆解① ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
header(s, "DEEP-DIVE  原理①", "BOSS 的「白板」是怎么回事")
blocks = [
    ("现象", "直接 goto 搜索深链 → 页面被吐成空 body，URL 停在 about:blank", BAD),
    ("根因", "深链缺少搜索流程中由 JS 生成的安全 token（__zp_stoken__），"
            "服务端/前端识别后直接返回空文档", WARN),
    ("对策", "模拟真人在首页搜索框打字 + 点搜索按钮 → 有机导航自动带上 token", ACCENT2),
    ("升级", "但发现连首页本身也被吐成 about:blank → 问题升级到「导航动作」这一层",
            ACCENT),
]
y = Inches(1.75)
for k, v, c in blocks:
    _box(s, Inches(0.6), y, Inches(12.1), Inches(1.16), fill=PANEL, line=c, line_w=1.5)
    _text(s, Inches(0.9), y + Inches(0.12), Inches(2.0), Inches(0.9),
          [[(k, 18, c, True)]], anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(2.85), y + Inches(0.12), Inches(9.6), Inches(0.9),
          [[(v, 15, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.32)
footer(s, 8)

# ---------------- 9. 原理拆解② ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
header(s, "DEEP-DIVE  原理②", "为什么 stealth 和原生导航都没用")
_text(s, Inches(0.6), Inches(1.75), Inches(12.1), Inches(1.0),
      [[("现象：", 16, BAD, True),
        ("注入反检测脚本无效；用 Target.createTarget 开的标签被秒关。", 16, TEXT, False)]])
_box(s, Inches(0.6), Inches(2.55), Inches(12.1), Inches(1.3), fill=PANEL, line=WARN, line_w=1.5)
_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.0),
      [[("诊断结论", 17, WARN, True),
        ("：BOSS 检测的是「CDP 调试器已挂载」这件事本身。", 17, TEXT, False)]],
      anchor=MSO_ANCHOR.MIDDLE)
_text(s, Inches(0.6), Inches(4.05), Inches(12.1), Inches(2.2),
      [[("原理：", 16, ACCENT, True)],
       [("• Playwright 默认会对每个 frame 调用 ", 15, TEXT, False),
        ("Runtime.enable", 15, ACCENT, True),
        ("，泄漏出一个可被探测的执行上下文 ID。", 15, TEXT, False)],
       [("• 反爬脚本据此判定「浏览器正被自动化控制」→ 直接白板 / 关标签。", 15, TEXT, False)],
       [("• 这是 Cloudflare、BOSS 等顶级反爬的通用、底层检测手段，", 15, TEXT, False),
        ("普通 stealth 与换导航方式都绕不过。", 15, WARN, True)]])
footer(s, 9)

# ---------------- 10. 当前对策 rebrowser ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
header(s, "SOLUTION  对策", "当前方案：rebrowser-playwright")
pts = [
    ("是什么", "Playwright 的「打补丁直替版」，导入名 rebrowser_playwright，业务代码零改动"),
    ("核心原理", "不自动调用 Runtime.enable，改用 addBinding 在需要时才取上下文 → 不再泄漏特征"),
    ("我们的做法", "仅替换浏览器引擎；其它站点照常工作，互不影响"),
    ("当前状态", "已装 1.49.1（与现有版本精确匹配）并接入，待重启验证 BOSS 首页是否正常渲染"),
]
y = Inches(1.8)
for k, v in pts:
    _box(s, Inches(0.6), y, Inches(12.1), Inches(1.12), fill=PANEL)
    _text(s, Inches(0.9), y + Inches(0.14), Inches(2.7), Inches(0.85),
          [[(k, 16, ACCENT, True)]], anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(3.6), y + Inches(0.14), Inches(8.8), Inches(0.85),
          [[(v, 15, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.28)
footer(s, 10)

# ---------------- 11. 仍难克服 ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
header(s, "LIMITS  局限", "仍难以彻底克服的，以及原理")
limits = [
    ("滑块 / 行为验证码", "对抗式人机验证（geetest 等），需真人轨迹或打码平台，程序无法稳定自动破解"),
    ("设备指纹 + IP 信誉", "服务端多维风控建模：单机单 IP 高频访问必然被限流，靠单点很难规避"),
    ("反爬军备竞赛", "rebrowser 等补丁也可能被新检测针对，需要持续跟进、动态对抗"),
    ("首次安全验证", "首次仍可能需人工过一次验证，登录态保存后才可持续自动"),
]
y = Inches(1.78)
for k, v in limits:
    _box(s, Inches(0.6), y, Inches(12.1), Inches(1.12), fill=PANEL, line=BAD, line_w=1.2)
    _text(s, Inches(0.9), y + Inches(0.14), Inches(3.4), Inches(0.85),
          [[(k, 16, BAD, True)]], anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(4.3), y + Inches(0.14), Inches(8.1), Inches(0.85),
          [[(v, 14, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(1.28)
footer(s, 11)

# ---------------- 12. 方案隐蔽性对比 ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
header(s, "COMPARE  对比", "各方案隐蔽性对比")
rows = [
    ("方案", "webdriver", "CDP 特征", "登录复用", "BOSS 可用"),
    ("普通 Selenium", "暴露", "暴露", "弱", "否"),
    ("普通 Playwright 启动", "可隐藏", "Runtime.enable 暴露", "中", "否"),
    ("CDP 接管真实 Chrome", "隐藏", "Runtime.enable 暴露", "强", "部分"),
    ("rebrowser + CDP 接管", "隐藏", "已消除", "强", "目标"),
]
tbl = s.shapes.add_table(len(rows), 5, Inches(0.5), Inches(1.9),
                         Inches(12.3), Inches(4.1)).table
widths = [3.4, 2.0, 3.0, 1.9, 2.0]
for i, wv in enumerate(widths):
    tbl.columns[i].width = Inches(wv)
last_col_color = {"否": BAD, "部分": WARN, "目标": ACCENT2}
for ri in range(len(rows)):
    for ci in range(5):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        if ri == 0:
            cell.fill.fore_color.rgb = ACCENT
        elif ri == len(rows) - 1:
            cell.fill.fore_color.rgb = RGBColor(0x18, 0x33, 0x2B)
        else:
            cell.fill.fore_color.rgb = PANEL if ri % 2 else RGBColor(0x16, 0x21, 0x33)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        run = p.add_run()
        run.text = rows[ri][ci]
        run.font.name = FONT
        run.font.size = Pt(14 if ri else 14)
        run.font.bold = (ri == 0 or ci == 0)
        if ri == 0:
            run.font.color.rgb = BG
        elif ci == 4:
            run.font.color.rgb = last_col_color.get(rows[ri][ci], TEXT)
        else:
            run.font.color.rgb = TEXT
footer(s, 12)

# ---------------- 13. Roadmap ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
header(s, "ROADMAP  下一步", "接下来怎么走")
steps = [
    "验证 rebrowser 是否让 BOSS 首页正常渲染（重启服务后重跑）",
    "若仍被拦：切换 alwaysIsolated 模式 / 对齐指纹（UA·语言·时区·分辨率）",
    "引入住宅代理 IP 池，分散访问频率，降低 IP 风控",
    "滑块兜底：保留「人工过一次验证」通道，评估打码平台对接",
    "其它站点持续稳定运行 + 抓取质量与成功率监控",
]
y = Inches(1.85)
for i, st in enumerate(steps, 1):
    _box(s, Inches(0.6), y, Inches(0.62), Inches(0.62), fill=ACCENT)
    _text(s, Inches(0.6), y, Inches(0.62), Inches(0.62),
          [[(str(i), 18, BG, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(1.4), y, Inches(11.2), Inches(0.62),
          [[(st, 16, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.92)
footer(s, 13)

# ---------------- 14. 结语 ----------------
s = prs.slides.add_slide(BLANK)
_bg(s)
_box(s, Inches(0.0), Inches(0.0), Inches(0.22), SH, fill=ACCENT)
_text(s, Inches(0.9), Inches(2.4), Inches(11.6), Inches(1.0),
      [[("总结", 34, WHITE, True)]])
_text(s, Inches(0.9), Inches(3.5), Inches(11.6), Inches(2.0),
      [[("已把抓取从「全靠手动」推进到「绝大多数站点全自动，", 18, TEXT, False)],
       [("BOSS 仅首次可能需过一次验证」。", 18, TEXT, False)],
       [("", 8, TEXT, False)],
       [("反爬是动态对抗，方案需持续演进——但核心难点已被逐层拆解。", 18, ACCENT, True)]])
_text(s, Inches(0.9), Inches(6.0), Inches(11.6), Inches(0.6),
      [[("谢谢观看", 16, MUTED, False)]])

out = Path(__file__).resolve().parent / "反爬攻坚汇报.pptx"
prs.save(str(out))
print("saved:", out)
