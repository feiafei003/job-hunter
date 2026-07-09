# -*- coding: utf-8 -*-
"""生成 Job Hunter Agent 介绍 PPT（精炼视觉版）。
运行：PYTHONPATH=./.pptx_lib python3 gen_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------- 设计系统 ----------------
BG     = RGBColor(0xF4, 0xF6, 0xFB)   # 页面浅底
INK    = RGBColor(0x0B, 0x12, 0x20)   # 标题近黑
BODY   = RGBColor(0x44, 0x52, 0x66)   # 正文
MUTED  = RGBColor(0x94, 0xA3, 0xB8)   # 次要
HAIR   = RGBColor(0xE4, 0xE9, 0xF2)   # 描边
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
NAVY   = RGBColor(0x0B, 0x12, 0x20)   # 封面底
NAVY2  = RGBColor(0x16, 0x21, 0x38)   # 封面卡

BLUE   = RGBColor(0x25, 0x63, 0xEB)
GREEN  = RGBColor(0x05, 0x96, 0x69)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
AMBER  = RGBColor(0xD9, 0x77, 0x06)
ROSE   = RGBColor(0xE1, 0x1D, 0x48)
CYAN   = RGBColor(0x0E, 0x74, 0x90)

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
PAGE = {"n": 0}


# ---------------- 基础工具 ----------------
def _set(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def soft_shadow(shape, alpha=22000, blur=0.10, dist=0.06):
    """给形状加柔和投影（python-pptx 无原生 API，直接写 XML）。"""
    spPr = shape._element.spPr
    for e in spPr.findall(qn('a:effectLst')):
        spPr.remove(e)
    eff = spPr.makeelement(qn('a:effectLst'), {})
    sh = eff.makeelement(qn('a:outerShdw'), {
        'blurRad': str(int(Inches(blur))),
        'dist': str(int(Inches(dist))),
        'dir': '5400000',
        'rotWithShape': '0',
    })
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '1E293B'})
    al = clr.makeelement(qn('a:alpha'), {'val': str(alpha)})
    clr.append(al)
    sh.append(clr)
    eff.append(sh)
    spPr.append(eff)


def rect(slide, x, y, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def rrect(slide, x, y, w, h, fill, line=None, radius=0.06, lw=1):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def oval(slide, x, y, d, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, color, bold, space) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = text
        _set(r, size, color, bold)
    return tb


def centered(slide, x, y, w, h, text, size, color, bold=True):
    textbox(slide, x, y, w, h, [(text, size, color, bold, 0)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def bg(slide):
    rect(slide, 0, 0, SW, SH, BG)


def header(slide, title, kicker=None):
    """浅色页眉：左侧竖条 + 标题 + kicker + 细分隔线。"""
    rrect(slide, Inches(0.6), Inches(0.52), Inches(0.16), Inches(0.66), BLUE, radius=0.5)
    textbox(slide, Inches(0.92), Inches(0.42), Inches(11.6), Inches(0.62),
            [(title, 27, INK, True, 0)], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        textbox(slide, Inches(0.94), Inches(1.12), Inches(11.6), Inches(0.34),
                [(kicker, 12.5, MUTED, False, 0)])
    rect(slide, Inches(0.6), Inches(1.56), Inches(12.13), Pt(1.2), HAIR)


def pageno(slide, n):
    textbox(slide, Inches(0.6), Inches(7.02), Inches(6), Inches(0.34),
            [("Job Hunter Agent", 9.5, MUTED, False, 0)])
    textbox(slide, Inches(11.7), Inches(7.02), Inches(1.05), Inches(0.34),
            [(f"{n:02d} / 13", 9.5, MUTED, False, 0)], align=PP_ALIGN.RIGHT)


def content_slide(title, kicker=None):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, title, kicker)
    PAGE["n"] += 1
    pageno(s, PAGE["n"])
    return s


def card(slide, x, y, w, h, title, desc, accent, num=None):
    """统一卡片：圆角 + 投影 + 彩色徽标/圆点 + 标题 + 细下划线 + 描述。"""
    c = rrect(slide, x, y, w, h, WHITE, line=HAIR, radius=0.07)
    soft_shadow(c)
    pad = Inches(0.32)
    if num is not None:
        b = oval(slide, x + pad, y + pad, Inches(0.42), accent)
        centered(slide, x + pad, y + pad, Inches(0.42), Inches(0.42), str(num), 15, WHITE)
        tx = x + pad + Inches(0.6)
        tw = w - pad - Inches(0.6) - Inches(0.2)
        textbox(slide, tx, y + pad - Inches(0.02), tw, Inches(0.5),
                [(title, 16.5, INK, True, 0)], anchor=MSO_ANCHOR.MIDDLE)
        ty = y + pad + Inches(0.62)
    else:
        oval(slide, x + pad, y + pad + Inches(0.04), Inches(0.16), accent)
        textbox(slide, x + pad + Inches(0.28), y + pad - Inches(0.06), w - pad - Inches(0.4), Inches(0.45),
                [(title, 16.5, INK, True, 0)])
        rect(slide, x + pad, y + pad + Inches(0.42), Inches(0.5), Pt(3), accent)
        ty = y + pad + Inches(0.62)
    tb = slide.shapes.add_textbox(x + pad, ty, w - pad - Inches(0.25), h - (ty - y) - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, d in enumerate(desc):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        r = p.add_run()
        r.text = d
        _set(r, 12.5, BODY, False)


def list_card(slide, x, y, w, h, title, items, accent, item_size=15, gap=12):
    """带标题与彩色圆点列表的白卡片。"""
    c = rrect(slide, x, y, w, h, WHITE, line=HAIR, radius=0.06)
    soft_shadow(c)
    pad = Inches(0.42)
    if title:
        oval(slide, x + pad, y + pad + Inches(0.05), Inches(0.16), accent)
        textbox(slide, x + pad + Inches(0.28), y + pad - Inches(0.05), w - pad * 2, Inches(0.45),
                [(title, 16.5, INK, True, 0)])
        rect(slide, x + pad, y + pad + Inches(0.44), Inches(0.55), Pt(3), accent)
        iy = y + pad + Inches(0.66)
    else:
        iy = y + pad
    tb = slide.shapes.add_textbox(x + pad, iy, w - pad * 2, h - (iy - y) - Inches(0.25))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r1 = p.add_run(); r1.text = "●  "; _set(r1, 9, accent, True)
        r2 = p.add_run(); r2.text = it; _set(r2, item_size, BODY, False)


# ================= 封面 =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, Inches(0.22), SH, BLUE)
textbox(s, Inches(1.0), Inches(2.1), Inches(11.5), Inches(1.3),
        [("AI 智能求职助手", 52, WHITE, True, 0)])
textbox(s, Inches(1.03), Inches(3.25), Inches(11.5), Inches(0.6),
        [("Job Hunter Agent", 22, RGBColor(0x8F, 0x9F, 0xBA), False, 0)])
rect(s, Inches(1.05), Inches(4.05), Inches(2.4), Pt(3), BLUE)
chips = [("多平台自动抓取", BLUE), ("AI 匹配分析", GREEN), ("个性化订阅推送", PURPLE)]
cx = Inches(1.03)
for txt, col in chips:
    w = Inches(0.5) + Inches(0.105) * len(txt)
    pill = rrect(s, cx, Inches(4.45), w, Inches(0.5), NAVY2, radius=0.5)
    oval(s, cx + Inches(0.2), Inches(4.45) + Inches(0.17), Inches(0.16), col)
    textbox(s, cx + Inches(0.42), Inches(4.45), w - Inches(0.4), Inches(0.5),
            [(txt, 13, RGBColor(0xCB, 0xD5, 0xE1), False, 0)], anchor=MSO_ANCHOR.MIDDLE)
    cx += w + Inches(0.25)
textbox(s, Inches(1.03), Inches(6.4), Inches(11.5), Inches(0.5),
        [("一个常驻服务器、端到端自动化的求职 Agent", 13.5, RGBColor(0x6B, 0x7A, 0x94), False, 0)])

# ================= 背景痛点 =================
s = content_slide("为什么做这个", "求职效率的三大痛点")
pains = [
    ("信息分散", ["职位散落在智联、BOSS、猎聘、前程无忧等多个平台", "每天反复刷、手动搜，重复且低效"]),
    ("筛选靠人肉", ["一条条点开看 JD，判断是否匹配自己", "海量职位里真正合适的很少，淘金成本高"]),
    ("容易错过", ["新职位时效性强，看晚了就被投满", "没有提醒机制，靠记忆和运气"]),
]
for i, (t, d) in enumerate(pains):
    card(s, Inches(0.6) + Inches(4.12) * i, Inches(1.85), Inches(3.85), Inches(3.0), t, d, ROSE)
band = rrect(s, Inches(0.6), Inches(5.2), Inches(12.13), Inches(1.1), RGBColor(0xEC, 0xF1, 0xFE), radius=0.12)
textbox(s, Inches(1.0), Inches(5.2), Inches(11.4), Inches(1.1),
        [("目标：把“找 → 筛 → 推”整条链路自动化，让合适的职位主动找到你。", 18, BLUE, True, 0)],
        anchor=MSO_ANCHOR.MIDDLE)

# ================= 产品概览 =================
s = content_slide("产品概览", "一句话：自动抓职位、AI 打分、按你的画像定时推荐")
caps = [
    ("多平台抓取", ["4 大主流招聘网站", "登录态保持 + 反爬", "（领英暂下线）"], BLUE),
    ("AI 匹配分析", ["结合个人画像打分 0–100", "亮点 / 风险 / 投递建议"], GREEN),
    ("智能订阅推送", ["按匹配度取最匹配 N 条", "定时邮件 + 去重"], PURPLE),
    ("对话式助手", ["简历分析 / 职位分析", "自然语言管理订阅"], AMBER),
]
for i, (t, d, c) in enumerate(caps):
    card(s, Inches(0.6) + Inches(3.07) * i, Inches(1.85), Inches(2.87), Inches(3.3), t, d, c)
band = rrect(s, Inches(0.6), Inches(5.5), Inches(12.13), Inches(0.95), WHITE, line=HAIR, radius=0.1)
soft_shadow(band)
textbox(s, Inches(1.0), Inches(5.5), Inches(11.5), Inches(0.95),
        [("技术栈：Python · FastAPI · SQLite · APScheduler · Playwright/CDP · LLM(Cursor gpt-5.5 / DeepSeek)",
          14.5, BODY, False, 0)], anchor=MSO_ANCHOR.MIDDLE)

# ================= 系统架构 =================
s = content_slide("系统架构", "一套常驻服务，模块解耦")
e1 = rrect(s, Inches(0.6), Inches(1.8), Inches(5.85), Inches(0.95), BLUE, radius=0.12); soft_shadow(e1)
centered(s, Inches(0.6), Inches(1.8), Inches(5.85), Inches(0.95), "用户端 Portal / 对话助手", 16, WHITE)
e2 = rrect(s, Inches(6.88), Inches(1.8), Inches(5.85), Inches(0.95), INK, radius=0.12); soft_shadow(e2)
centered(s, Inches(6.88), Inches(1.8), Inches(5.85), Inches(0.95), "管理端 Admin（可分端口）", 16, WHITE)
core = rrect(s, Inches(0.6), Inches(3.0), Inches(12.13), Inches(0.85), RGBColor(0x33, 0x41, 0x55), radius=0.1)
soft_shadow(core)
centered(s, Inches(0.6), Inches(3.0), Inches(12.13), Inches(0.85),
         "FastAPI 应用核心  ·  鉴权  ·  REST API", 15.5, WHITE)
mods = [
    ("APScheduler", ["夜间预热 + 时段", "预抓取 + 推送"], BLUE),
    ("抓取插件 CDP", ["4 站点 scraper", "列表/详情解耦"], GREEN),
    ("AI 分析器 LLM", ["匹配度打分", "建议生成"], PURPLE),
    ("订阅 / 推送", ["匹配 → 补详情", "排序 → 发信"], AMBER),
]
for i, (t, d, c) in enumerate(mods):
    cxp = Inches(0.6) + Inches(3.07) * i
    mc = rrect(s, cxp, Inches(4.1), Inches(2.87), Inches(1.7), WHITE, line=HAIR, radius=0.1)
    soft_shadow(mc)
    rect(s, cxp + Inches(0.32), Inches(4.35), Inches(0.45), Pt(3.5), c)
    textbox(s, cxp + Inches(0.32), Inches(4.45), Inches(2.4), Inches(0.45),
            [(t, 14.5, INK, True, 0)])
    tb = s.shapes.add_textbox(cxp + Inches(0.32), Inches(4.95), Inches(2.45), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True
    for j, line in enumerate(d):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run(); r.text = line; _set(r, 11.5, BODY, False)
foot = rrect(s, Inches(0.6), Inches(6.05), Inches(12.13), Inches(0.75), RGBColor(0xEC, 0xF1, 0xFE), radius=0.12)
textbox(s, Inches(0.6), Inches(6.05), Inches(12.13), Inches(0.75),
        [("存储 SQLite（职位 / 分析 / 订阅 / 投递）  ·  浏览器经中国代理出网，破解地域限制",
          13, RGBColor(0x33, 0x49, 0x7A), False, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ================= 功能一 抓取 =================
s = content_slide("功能一 · 多平台职位抓取", "一次配置，周期自动执行")
list_card(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(4.85), "", [
    "覆盖智联招聘、BOSS直聘、猎聘、前程无忧（领英暂下线）",
    "插件式架构：新增站点只需加一个 scraper",
    "持久化登录态：手动登录一次，后续免重复登录",
    "真实可视化浏览器：扫码、风控可远程 VNC 接管",
    "反爬补丁版内核 + 随机延迟，降低风控",
    "按 URL 指纹去重入库，避免重复",
], BLUE, item_size=15.5, gap=15)
card(s, Inches(6.85), Inches(1.85), Inches(5.88), Inches(2.3), "登录态保持",
     ["基于 Playwright 持久化上下文，登录信息本地保存",
      "遇验证码 / 滑块可在窗口手动处理，程序自动继续"], GREEN)
card(s, Inches(6.85), Inches(4.4), Inches(5.88), Inches(2.3), "可配置搜索任务",
     ["站点 + 关键词 + 城市 + 频率，网页即可增删",
      "支持立即执行，也可交给调度器周期跑"], BLUE)

# ================= 分层职位池 =================
s = content_slide("分层职位池与按需详情", "抓取(列表) 与 详情+AI 解耦，又快又省")
layers = [
    ("夜间基础预热", ["每天 03:00 自动跑", "8 城 × 全角色，只抓列表灌入中央池", "新用户订阅后即刻有内容"], BLUE),
    ("按需详情 + AI", ["点「获取详情 & AI 分析」才触发", "抓该条完整 JD + 按本人画像打分", "全屏展示，已抓过走缓存"], GREEN),
    ("订阅发送补抓", ["先按列表规则匹配候选", "对入选职位逐条补详情 + 分析", "按匹配度取 Top N 再发信"], PURPLE),
]
for i, (t, d, c) in enumerate(layers):
    card(s, Inches(0.6) + Inches(4.12) * i, Inches(1.85), Inches(3.85), Inches(3.5), t, d, c, num=i + 1)
band = rrect(s, Inches(0.6), Inches(5.7), Inches(12.13), Inches(1.0), RGBColor(0xEC, 0xF1, 0xFE), radius=0.12)
textbox(s, Inches(1.0), Inches(5.7), Inches(11.4), Inches(1.0),
        [("列表态足够“今日匹配”浏览；完整 JD 与打分只在用户感兴趣或要发邮件时按需产生，省算力、出结果更快。",
          15.5, BLUE, True, 0)], anchor=MSO_ANCHOR.MIDDLE)

# ================= 功能二 AI 分析 =================
s = content_slide("功能二 · AI 匹配分析", "把“看 JD”这件事交给 AI")
list_card(s, Inches(0.6), Inches(1.85), Inches(6.0), Inches(4.85), "", [
    "结构化「求职画像」：经验 / 技能 / 目标岗位 / 期望薪资",
    "LLM 结合画像对每个职位输出 0–100 匹配度",
    "同时给出一句话结论、投递建议、补强技能、简历改进",
    "分析结果带画像指纹缓存，避免重复计算",
    "LLM 后端可切换：Cursor(gpt-5.5) / DeepSeek",
], GREEN, item_size=15.5, gap=16)
demo = rrect(s, Inches(6.85), Inches(1.85), Inches(5.88), Inches(4.85), WHITE, line=HAIR, radius=0.07)
soft_shadow(demo)
top = rrect(s, Inches(6.85), Inches(1.85), Inches(5.88), Inches(0.8), PURPLE, radius=0.07)
centered(s, Inches(6.85), Inches(1.85), Inches(5.88), Inches(0.8), "AI 分析卡（示意）", 15, WHITE)
oval(s, Inches(7.2), Inches(2.95), Inches(1.2), GREEN)
centered(s, Inches(7.2), Inches(2.95), Inches(1.2), Inches(1.2), "88", 30, WHITE)
textbox(s, Inches(8.7), Inches(2.95), Inches(3.8), Inches(1.2),
        [("匹配度 88", 18, INK, True, 4), ("与你的 C++ 后端背景高度契合", 12.5, BODY, False, 0)],
        anchor=MSO_ANCHOR.MIDDLE)
textbox(s, Inches(7.2), Inches(4.4), Inches(5.2), Inches(2.1),
        [
            ("建议：强调高并发与 Linux 网络编程经验。", 13, BODY, False, 9),
            ("需补强：Kubernetes 实战经验。", 13, BODY, False, 9),
            ("简历改进：把 QPS / 延迟等量化指标写出来。", 13, BODY, False, 0),
        ])

# ================= 功能三 订阅 =================
s = content_slide("功能三 · 智能订阅与邮件推送", "重点：只推“最匹配的 N 条”")
list_card(s, Inches(0.6), Inches(1.85), Inches(6.05), Inches(4.85), "", [
    "过滤条件：关键词 / 站点 / 城市 / 薪资 / 工作类型",
    "发送时段：每天 10:00、21:00、工作日、每周一（cron）",
    "发送前自动补抓详情，保证内容新鲜",
    "已推送去重：同一职位不重复打扰",
    "邮件内嵌 AI 分析与匹配度徽标，一眼看懂",
], BLUE, item_size=15, gap=14)
panel = rrect(s, Inches(6.9), Inches(1.85), Inches(5.83), Inches(4.85), WHITE, line=HAIR, radius=0.07)
soft_shadow(panel)
top = rrect(s, Inches(6.9), Inches(1.85), Inches(5.83), Inches(0.8), GREEN, radius=0.07)
centered(s, Inches(6.9), Inches(1.85), Inches(5.83), Inches(0.8), "本次优化的发送设置", 15.5, WHITE)
tb = s.shapes.add_textbox(Inches(7.25), Inches(2.9), Inches(5.15), Inches(3.6))
tf = tb.text_frame; tf.word_wrap = True
opts = [
    "每次最多 N 条：用户可自定义（1–50）",
    "按 AI 匹配度从高到低取 Top N（真·最匹配）",
    "最低匹配度门槛：低于不推",
    "无匹配也可选发提醒，或静默不打扰",
    "邮件标题与卡片显示匹配分数",
]
for i, it in enumerate(opts):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(13)
    r1 = p.add_run(); r1.text = "●  "; _set(r1, 9, GREEN, True)
    r2 = p.add_run(); r2.text = it; _set(r2, 14.5, INK, False)

# ================= 功能四 助手 =================
s = content_slide("功能四 · 对话式助手", "用聊天的方式完成操作")
cs = [
    ("简历分析", ["上传 / 粘贴简历", "AI 给出优劣势与改进建议"], BLUE),
    ("职位分析", ["贴一条 JD 或选已抓职位", "结合画像即时打分与解读"], GREEN),
    ("订阅管理", ["“帮我订上海 C++，每天 5 条”", "自然语言增删改订阅"], PURPLE),
]
for i, (t, d, c) in enumerate(cs):
    card(s, Inches(0.6) + Inches(4.12) * i, Inches(1.95), Inches(3.85), Inches(3.2), t, d, c)
band = rrect(s, Inches(0.6), Inches(5.45), Inches(12.13), Inches(0.95), WHITE, line=HAIR, radius=0.1)
soft_shadow(band)
textbox(s, Inches(1.0), Inches(5.45), Inches(11.4), Inches(0.95),
        [("助手内置工具调用 + 关键词预设（6 大类约 30 个角色），新手也能快速配置。", 15.5, BODY, False, 0)],
        anchor=MSO_ANCHOR.MIDDLE)

# ================= 用户门户与上手引导 =================
s = content_slide("用户门户与上手引导", "账号体系 + 画像沉淀，开箱即用")
list_card(s, Inches(0.6), Inches(1.85), Inches(6.05), Inches(4.85), "", [
    "统一账号：注册 / 登录，个人画像长期沉淀",
    "注册即引导完善画像，匹配从第一天就更准",
    "首推上传简历，AI 自动解析填充，可“稍后再说”",
    "「我的画像」展示优先：默认只读，一键修改后保存",
    "画像为空时今日匹配顶部软提示，引导去完善",
    "匿名订阅页下线，统一走账号（退订链接仍有效）",
], PURPLE, item_size=15, gap=13)
card(s, Inches(6.9), Inches(1.85), Inches(5.83), Inches(2.3), "注册引导（新用户）",
     ["注册 → 完善画像 → 开始匹配的顺滑动线",
      "上传简历一键填充，降低上手成本"], GREEN)
card(s, Inches(6.9), Inches(4.4), Inches(5.83), Inches(2.3), "画像展示优先",
     ["有画像默认只读展示，点「修改」再编辑",
      "订阅未单独填画像时，默认复用「我的画像」"], BLUE)

# ================= 技术亮点 =================
s = content_slide("技术亮点", "难点与解决方案")
ts = [
    ("可视化浏览器自动化", ["headful + VNC，真实窗口扫码登录", "CDP(9222) 远程接管处理风控"], BLUE),
    ("地域限制破解", ["招聘站 / LLM 对境外 IP 限流", "经中国代理出网，稳定访问"], GREEN),
    ("LLM 工程化", ["从不稳定 SDK 改为直接调 CLI", "区域可用模型自动适配"], PURPLE),
    ("分层抓取调度", ["夜间预热 + 时段两段式", "列表/详情解耦，按需才进详情"], AMBER),
]
for i, (t, d, c) in enumerate(ts):
    cxp = Inches(0.6) + Inches(6.18) * (i % 2)
    cyp = Inches(1.85) + Inches(2.45) * (i // 2)
    card(s, cxp, cyp, Inches(5.95), Inches(2.2), t, d, c)

# ================= 工程与部署 =================
s = content_slide("工程与部署", "面向长期常驻运行")
list_card(s, Inches(0.6), Inches(1.85), Inches(12.13), Inches(4.85), "", [
    "用户端 / 管理端按 APP_ROLE 拆分，可同端口或分端口部署",
    "对外门户固定监听 0.0.0.0，setsid 脱离会话常驻、不随终端退出",
    "run_server.sh 守护进程：崩溃自动退避重启，纯用户态、无需 root",
    "Linux 服务器常驻；无桌面也能跑，扫码时用 VNC 看真实窗口",
    "SQLite 轻量存储 + 启动时自动轻量迁移，升级不丢数据",
    "数据与登录态本地化，隐私可控",
], CYAN, item_size=17, gap=16)

# ================= 价值 =================
s = content_slide("价值与效果")
vs = [
    ("省时间", ["不再手动刷多个网站", "新职位主动送上门"], GREEN),
    ("更精准", ["AI 按画像打分排序", "只看最匹配的少数几条"], BLUE),
    ("不错过", ["定时推送 + 预抓取", "时效性强的机会及时触达"], PURPLE),
]
for i, (t, d, c) in enumerate(vs):
    card(s, Inches(0.6) + Inches(4.12) * i, Inches(1.95), Inches(3.85), Inches(3.0), t, d, c)
band = rrect(s, Inches(0.6), Inches(5.3), Inches(12.13), Inches(1.05), RGBColor(0xEC, 0xF1, 0xFE), radius=0.12)
textbox(s, Inches(1.0), Inches(5.3), Inches(11.4), Inches(1.05),
        [("从“人找工作”到“工作找人”——把重复劳动交给 Agent，把判断留给自己。", 18, BLUE, True, 0)],
        anchor=MSO_ANCHOR.MIDDLE)

# ================= 规划 =================
s = content_slide("后续规划")
list_card(s, Inches(0.6), Inches(1.85), Inches(12.13), Inches(4.85), "", [
    "更多渠道推送：微信 / Server酱 / Webhook",
    "更灵活的发送时段与去重时间窗",
    "一键投递 / 投递追踪",
    "多账号、多画像的并行管理",
    "面试问题预测与准备建议",
], AMBER, item_size=17, gap=16)

# ================= 谢谢 =================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, Inches(0.22), SH, BLUE)
centered(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(1.2), "谢谢观看", 46, WHITE)
rect(s, Inches(6.16), Inches(4.15), Inches(1.0), Pt(3), BLUE)
centered(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(0.7),
         "AI 智能求职助手 · Job Hunter Agent", 17, RGBColor(0x8F, 0x9F, 0xBA), False)

prs.save("Job_Hunter_Agent.pptx")
print("saved", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
